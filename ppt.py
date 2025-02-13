from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation
ppt = Presentation()

# Define slide titles and content
slides_content = [
    ("Smart Car Parking System", "Simulation and Cloud Integration\n\nB.Tech IT - ABC Company Collaboration"),
    ("Introduction", "The project simulates a Smart Car Parking System using Wokwi and VS Code.\nGoogle Sheets is used for cloud storage.\nReal-time monitoring via an HTML page on GitHub Pages."),
    ("System Features", "1. User Interaction via Buttons and Sensors\n2. Real-Time Monitoring\n3. Cloud Integration with Google Sheets\n4. Scalability Demonstration"),
    ("User Interaction via Buttons", "- Entry Confirmation (Button 1): Opens gate, reduces available spaces.\n- Exit Confirmation (Button 2): Opens gate, increases available spaces.\n- Manual Update (Button 3): Adjusts spaces for errors or maintenance."),
    ("Real-Time Monitoring", "- Display available spaces on an LCD and HTML page.\n- Users can monitor but not modify data."),
    ("Cloud Integration", "- Google Sheets API manages parking space availability.\n- Data is sent and retrieved in real-time."),
    ("Scalability Demonstration", "- System runs on three different computers simultaneously.\n- Ensures multiple users can access real-time data without conflicts."),
    ("Implementation Tools", "1. Wokwi: Hardware simulation\n2. VS Code: Coding and debugging\n3. Google Sheets API: Cloud storage\n4. HTML/JavaScript: Monitoring page"),
    ("Tasks", "1. Simulate the system using Wokwi.\n2. Integrate Google Sheets for real-time data.\n3. Develop an HTML page for monitoring.\n4. Demonstrate scalability on multiple computers."),
    ("Evaluation Criteria", "1. Correct simulation functionality.\n2. Accurate Google Sheets integration.\n3. Usability of HTML monitoring page.\n4. Successful scalability demonstration."),
    ("System Design & Architecture", "- Block diagram\n- Communication flow\n- Software & hardware interaction\n- Data flow from sensors to cloud & HTML page"),
    ("Implementation Steps", "- Coding & API integration\n- Setting up Wokwi\n- Configuring Google Sheets API\n- Designing HTML monitoring page"),
    ("Challenges & Solutions", "- Real-time data synchronization issues\n- Network failures\n- Optimizing API requests"),
    ("Future Enhancements", "- AI-based parking space prediction\n- IoT sensors for vehicle detection\n- RFID access control for security"),
    ("Conclusion", "The system integrates embedded functionalities with cloud storage.\nSuccessful implementation demonstrates feasibility for real-world applications.")
]

# Add slides to the presentation
for title, content in slides_content:
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and Content layout
    slide.shapes.title.text = title
    slide.shapes.placeholders[1].text = content

# Save the PowerPoint file
ppt_file = "/mnt/data/Smart_Car_Parking_System.pptx"
ppt.save(ppt_file)

ppt_file
